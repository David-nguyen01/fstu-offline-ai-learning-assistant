from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from .text_utils import normalize_text


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


class DocumentLoadError(RuntimeError):
    pass


@dataclass(frozen=True)
class DocumentPage:
    text: str
    page: int | None
    source: str


def load_document(path: Path) -> list[DocumentPage]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise DocumentLoadError(f"Unsupported file type: {suffix}")
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    return _load_text(path)


def _load_pdf(path: Path) -> list[DocumentPage]:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError as exc:
            raise DocumentLoadError("Install pypdf to read PDF files.") from exc

    reader = PdfReader(str(path))
    pages: list[DocumentPage] = []
    for index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text:
            pages.append(DocumentPage(text=text, page=index, source=path.name))
    if not pages:
        raise DocumentLoadError("No readable text found in this PDF.")
    return pages


def _load_docx(path: Path) -> list[DocumentPage]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentLoadError("Install python-docx to read DOCX files.") from exc

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [" ".join(cell.text.split()) for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    text = normalize_text("\n".join(paragraphs))
    if not text:
        raise DocumentLoadError("No readable text found in this DOCX.")
    return [DocumentPage(text=text, page=None, source=path.name)]


def _load_pptx(path: Path) -> list[DocumentPage]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentLoadError("Install python-pptx to read PPTX files.") from exc

    presentation = Presentation(str(path))
    pages: list[DocumentPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        text = normalize_text("\n".join(texts))
        if text:
            pages.append(DocumentPage(text=text, page=index, source=path.name))
    if not pages:
        raise DocumentLoadError("No readable text found in this PPTX.")
    return pages


def _load_text(path: Path) -> list[DocumentPage]:
    raw = path.read_bytes()
    text: str | None = None
    for encoding in ["utf-8-sig", "utf-8", "utf-16", "cp1258", "latin-1"]:
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    text = normalize_text(text or raw.decode("utf-8", errors="replace"))
    if not text:
        raise DocumentLoadError("No readable text found in this text file.")
    return [DocumentPage(text=text, page=None, source=path.name)]
